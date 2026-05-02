from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Color palette ──────────────────────────────────────────────
C_BG_DARK   = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
C_BG2       = RGBColor(0x0F, 0x2A, 0x1A)   # deep green-navy
C_GREEN     = RGBColor(0x4A, 0xD2, 0x95)   # mint green
C_BLUE      = RGBColor(0x38, 0xBD, 0xF8)   # sky blue
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY      = RGBColor(0xB0, 0xB8, 0xC8)
C_CARD      = RGBColor(0x1A, 0x2A, 0x3A)
C_ACCENT2   = RGBColor(0x81, 0x8C, 0xF8)   # indigo
C_GOLD      = RGBColor(0xF5, 0xC5, 0x42)

blank_layout = prs.slide_layouts[6]  # completely blank

# ── Helpers ────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line_color=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.width = line_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_w if line_w else Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size=18, bold=False, color=None, align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color if color else C_WHITE
    return txBox

def bg(slide, color=C_BG_DARK, color2=None):
    """Full-slide background"""
    add_rect(slide, 0, 0, 13.33, 7.5, fill=color)
    if color2:
        # second half gradient approximation
        add_rect(slide, 6.5, 0, 6.83, 7.5, fill=color2)

def accent_bar(slide, l=0.7, t=2.1, w=0.55, h=0.07):
    add_rect(slide, l, t, w, h, fill=C_GREEN)

def section_label(slide, text, l=0.7, t=0.4):
    add_text(slide, text, l, t, 11, 0.4, size=10, color=RGBColor(0x4A,0xD2,0x95), bold=True)

def slide_title(slide, title, l=0.7, t=0.75, w=11.9, size=36):
    add_text(slide, title, l, t, w, 0.8, size=size, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)

def card(slide, l, t, w, h, icon, title, body, icon_size=22, title_size=14, body_size=11):
    add_rect(slide, l, t, w, h, fill=C_CARD, line_color=RGBColor(0x2A,0x3A,0x4A), line_w=Pt(1))
    add_text(slide, icon, l+0.18, t+0.18, 0.6, 0.5, size=icon_size)
    add_text(slide, title, l+0.18, t+0.75, w-0.36, 0.35, size=title_size, bold=True, color=C_GREEN)
    add_text(slide, body, l+0.18, t+1.12, w-0.36, h-1.3, size=body_size, color=C_GRAY, wrap=True)

def screenshot_placeholder(slide, l, t, w, h, label):
    add_rect(slide, l, t, w, h, fill=RGBColor(0x12,0x22,0x1A), line_color=C_GREEN, line_w=Pt(1.5))
    add_text(slide, "📱", l, t+h/2-0.5, w, 0.6, size=28, align=PP_ALIGN.CENTER)
    add_text(slide, "[ Insérez votre capture ]", l, t+h/2+0.1, w, 0.4, size=9, color=C_GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, label, l, t+h-0.4, w, 0.35, size=9, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)

def deco_circle(slide, l, t, size=2.5, alpha_color=RGBColor(0x15,0x2A,0x20)):
    shape = slide.shapes.add_shape(9, Inches(l), Inches(t), Inches(size), Inches(size))
    shape.fill.background()
    shape.line.color.rgb = alpha_color
    shape.line.width = Pt(1.2)

# ══════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
bg(sl, C_BG_DARK, RGBColor(0x0D,0x22,0x16))
deco_circle(sl, 9.5, -1.2, 4.5)
deco_circle(sl, -1.5, 4.5, 4.0)
deco_circle(sl, 10.5, 4.5, 3.0)

# Green top accent strip
add_rect(sl, 0, 0, 0.12, 7.5, fill=C_GREEN)

# Tag pill
add_rect(sl, 0.9, 1.0, 3.2, 0.38, fill=RGBColor(0x0A,0x30,0x20), line_color=C_GREEN, line_w=Pt(1))
add_text(sl, "PROJET MOBILE  •  FLUTTER", 0.9, 1.0, 3.2, 0.38, size=9, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)

# App name
add_text(sl, "Curan", 0.9, 1.55, 8, 1.5, size=72, bold=True, color=C_WHITE)
# Arabic subtitle
add_text(sl, "قُرْآن", 0.9, 3.0, 5, 0.9, size=36, color=C_GREEN)
# Description
add_text(sl, "Application mobile sécurisée pour écouter le Coran\navec authentification biométrique, Firebase & horaires de prière.", 0.9, 3.95, 7.5, 0.9, size=14, color=C_GRAY)

# Author box
add_rect(sl, 0.9, 5.1, 5.5, 1.7, fill=RGBColor(0x0A,0x18,0x12), line_color=RGBColor(0x2A,0x4A,0x35), line_w=Pt(1))
add_text(sl, "PRÉSENTÉ PAR", 1.1, 5.25, 5, 0.3, size=8, bold=True, color=RGBColor(0x60,0x90,0x78))
add_text(sl, "[ Votre Nom Complet ]", 1.1, 5.55, 5, 0.45, size=20, bold=True, color=C_GREEN)
add_text(sl, "[ Filière  •  Promotion  •  Établissement ]", 1.1, 6.0, 5, 0.3, size=10, color=C_GRAY)
add_text(sl, "[ Date ]", 1.1, 6.35, 5, 0.3, size=9, color=RGBColor(0x50,0x70,0x60))

# Right decoration — moon
add_text(sl, "🌙", 10.0, 2.0, 2.5, 2.5, size=100, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SLIDE 2 — SOMMAIRE
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
bg(sl, C_BG_DARK, RGBColor(0x0D,0x1E,0x16))
add_rect(sl, 0, 0, 0.12, 7.5, fill=C_GREEN)
deco_circle(sl, 9.5, -1, 4)

section_label(sl, "PLAN DE PRÉSENTATION")
slide_title(sl, "Sommaire")
accent_bar(sl, l=0.7, t=1.65)

items = [
    ("01", "🎯", "Présentation du Projet", "Contexte, objectifs et motivation"),
    ("02", "✨", "Fonctionnalités Clés", "Les 8 features principales de l'application"),
    ("03", "🏗️", "Architecture Technique", "Stack, patterns et organisation du code"),
    ("04", "📱", "Démonstration", "Aperçu des écrans et interfaces"),
    ("05", "💪", "Défis & Solutions", "Problèmes rencontrés et solutions adoptées"),
    ("06", "🔭", "Conclusion & Perspectives", "Bilan et améliorations futures"),
]

for i, (num, icon, title, sub) in enumerate(items):
    col = i % 3
    row = i // 3
    lx = 0.7 + col * 4.2
    ty = 2.0 + row * 2.2
    add_rect(sl, lx, ty, 3.9, 2.0, fill=C_CARD, line_color=RGBColor(0x2A,0x3A,0x4A), line_w=Pt(1))
    add_text(sl, num, lx+0.18, ty+0.15, 0.8, 0.45, size=11, bold=True, color=C_GREEN)
    add_text(sl, icon, lx+3.1, ty+0.1, 0.65, 0.5, size=18, align=PP_ALIGN.RIGHT)
    add_text(sl, title, lx+0.18, ty+0.62, 3.5, 0.42, size=13, bold=True, color=C_WHITE)
    add_text(sl, sub, lx+0.18, ty+1.08, 3.5, 0.7, size=10, color=C_GRAY, wrap=True)

# ══════════════════════════════════════════════════════════════
# SLIDE 3 — PRÉSENTATION DU PROJET
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
bg(sl, C_BG_DARK, RGBColor(0x0D,0x20,0x16))
add_rect(sl, 0, 0, 0.12, 7.5, fill=C_GREEN)
deco_circle(sl, 9.8, -0.8, 3.5)

section_label(sl, "01 — PRÉSENTATION")
slide_title(sl, "Qu'est-ce que Curan ?")
accent_bar(sl, l=0.7, t=1.65)

card(sl, 0.7,  1.85, 3.9, 2.2, "📖", "Contexte",
     "Application Flutter dédiée à l'écoute du Saint Coran.\nDisponible sur Android, iOS et Web.")
card(sl, 4.72, 1.85, 3.9, 2.2, "🎯", "Objectif",
     "Offrir une expérience sécurisée et personnalisée pour la récitation et l'écoute du Coran.")
card(sl, 8.74, 1.85, 3.9, 2.2, "💡", "Motivation",
     "Combiner spiritualité et technologie : biométrie, cloud, géolocalisation et statistiques.")

# Tech badges row
badges = ["Flutter 3.x", "Firebase", "Dart", "Multi-plateforme", "Open Source"]
bx = 0.7
for b in badges:
    bw = len(b)*0.12 + 1.0
    add_rect(sl, bx, 4.3, bw, 0.42, fill=RGBColor(0x0A,0x28,0x1A), line_color=C_GREEN, line_w=Pt(1))
    add_text(sl, b, bx+0.08, 4.3, bw-0.1, 0.42, size=11, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)
    bx += bw + 0.2

add_text(sl, "قُرْآن", 9.5, 4.0, 3.5, 3.0, size=72, color=RGBColor(0x15,0x35,0x25), align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SLIDE 4 — FONCTIONNALITÉS
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
bg(sl, C_BG_DARK, RGBColor(0x0F,0x22,0x18))
add_rect(sl, 0, 0, 0.12, 7.5, fill=C_GREEN)
deco_circle(sl, 9.5, 4, 4)

section_label(sl, "02 — FONCTIONNALITÉS")
slide_title(sl, "Fonctionnalités Clés")
accent_bar(sl, l=0.7, t=1.65)

features = [
    ("🔒", "Authentification Firebase",        "Connexion, inscription et mot de passe oublié"),
    ("🪪", "Authentification Biométrique",      "Empreinte digitale et Face ID via local_auth"),
    ("🎵", "Lecteur Audio Coran",               "Contrôle complet : lecture, pause, navigation"),
    ("❤️", "Gestion des Favoris",              "Sauvegarde des sourates préférées en cloud"),
    ("🕌", "Horaires de Prière",               "Calcul précis par GPS et géolocalisation"),
    ("📊", "Dashboard & Statistiques",         "Graphiques d'écoute avec fl_chart"),
    ("☁️", "Synchronisation Cloud",            "Firestore temps réel pour données utilisateur"),
    ("⚙️", "Paramètres Personnalisables",      "Thème, récitateur, notifications"),
]

for i, (icon, title, body) in enumerate(features):
    col = i % 4
    row = i // 4
    lx = 0.7 + col * 3.15
    ty = 1.85 + row * 2.55
    add_rect(sl, lx, ty, 3.0, 2.3, fill=C_CARD, line_color=RGBColor(0x2A,0x3A,0x4A), line_w=Pt(1))
    add_text(sl, icon, lx+0.15, ty+0.15, 0.7, 0.55, size=20)
    add_text(sl, title, lx+0.15, ty+0.75, 2.7, 0.45, size=11, bold=True, color=C_GREEN, wrap=True)
    add_text(sl, body,  lx+0.15, ty+1.22, 2.7, 0.9,  size=9,  color=C_GRAY,  wrap=True)

# ══════════════════════════════════════════════════════════════
# SLIDE 5 — ARCHITECTURE
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
bg(sl, C_BG_DARK, RGBColor(0x0D,0x1A,0x2A))
add_rect(sl, 0, 0, 0.12, 7.5, fill=C_GREEN)
deco_circle(sl, 10, -1, 4)

section_label(sl, "03 — ARCHITECTURE")
slide_title(sl, "Stack Technique & Architecture")
accent_bar(sl, l=0.7, t=1.65)

# Architecture layers
layers = [
    ("📱", "UI Layer",   "Screens\nWidgets\nMaterial"),
    ("⚙️", "State",     "Provider\nAuthProvider\nAudioProvider"),
    ("🔧", "Services",  "AudioService\nAuthService\nBiometricService"),
    ("☁️", "Backend",  "Firebase Auth\nFirestore\nREST API"),
]
arrow_positions = [3.3, 5.55, 7.8]

for i, (icon, lname, items_text) in enumerate(layers):
    lx = 0.7 + i * 2.9
    ty = 2.1
    add_rect(sl, lx, ty, 2.6, 2.8, fill=C_CARD, line_color=RGBColor(0x2A,0x3A,0x4A), line_w=Pt(1))
    add_text(sl, icon,       lx+0.1, ty+0.15, 2.4, 0.55, size=22, align=PP_ALIGN.CENTER)
    add_text(sl, lname,      lx+0.1, ty+0.75, 2.4, 0.38, size=12, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)
    add_text(sl, items_text, lx+0.1, ty+1.18, 2.4, 1.4,  size=9,  color=C_GRAY, align=PP_ALIGN.CENTER, wrap=True)

for ax in arrow_positions:
    add_text(sl, "→", ax, 3.1, 0.5, 0.6, size=18, color=C_GREEN, align=PP_ALIGN.CENTER)

# Packages row
pkgs = ["just_audio", "fl_chart", "local_auth", "geolocator", "shared_prefs", "http", "intl"]
add_text(sl, "Packages principaux :", 0.7, 5.15, 3.5, 0.4, size=10, bold=True, color=C_GRAY)
bx = 0.7
for p in pkgs:
    bw = len(p)*0.11 + 0.8
    add_rect(sl, bx, 5.6, bw, 0.38, fill=RGBColor(0x0D,0x20,0x18), line_color=RGBColor(0x2A,0x4A,0x35), line_w=Pt(1))
    add_text(sl, p, bx+0.05, 5.6, bw-0.1, 0.38, size=9, color=C_GREEN, align=PP_ALIGN.CENTER)
    bx += bw + 0.18
    if bx > 12.5: break

# ══════════════════════════════════════════════════════════════
# SLIDE 6 — SCREENSHOTS AUTH
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
bg(sl, C_BG_DARK, RGBColor(0x0D,0x20,0x16))
add_rect(sl, 0, 0, 0.12, 7.5, fill=C_GREEN)

section_label(sl, "04 — DÉMONSTRATION")
slide_title(sl, "Authentification & Sécurité")
accent_bar(sl, l=0.7, t=1.65)

screens_auth = ["Splash Screen", "Connexion", "Inscription", "Biométrie"]
for i, name in enumerate(screens_auth):
    lx = 0.7 + i * 3.15
    screenshot_placeholder(sl, lx, 1.85, 2.85, 5.0, name)

# ══════════════════════════════════════════════════════════════
# SLIDE 7 — SCREENSHOTS PLAYER
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
bg(sl, C_BG_DARK, RGBColor(0x0D,0x1E,0x18))
add_rect(sl, 0, 0, 0.12, 7.5, fill=C_GREEN)

section_label(sl, "04 — DÉMONSTRATION")
slide_title(sl, "Lecteur Audio & Dashboard")
accent_bar(sl, l=0.7, t=1.65)

screens_player = ["Dashboard", "Lecteur Audio", "Favoris", "Horaires Prière"]
for i, name in enumerate(screens_player):
    lx = 0.7 + i * 3.15
    screenshot_placeholder(sl, lx, 1.85, 2.85, 5.0, name)

# ══════════════════════════════════════════════════════════════
# SLIDE 8 — DÉFIS & SOLUTIONS
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
bg(sl, C_BG_DARK, RGBColor(0x0D,0x1A,0x2A))
add_rect(sl, 0, 0, 0.12, 7.5, fill=C_GREEN)
deco_circle(sl, 10, 4, 4)

section_label(sl, "05 — DÉFIS & SOLUTIONS")
slide_title(sl, "Défis Techniques Rencontrés")
accent_bar(sl, l=0.7, t=1.65)

challenges = [
    ("🔐", "Intégration Biométrique",
     "Compatibilité multi-plateforme avec local_auth.\nGestion des fallbacks et permissions."),
    ("🎵", "Lecteur Audio Persistant",
     "Maintien de la lecture entre les navigations via PersistentPlayerWrapper."),
    ("📡", "Synchronisation Temps Réel",
     "Sync favoris & stats avec Firestore sans conflits ni perte de données."),
    ("🕌", "Horaires de Prière GPS",
     "Calcul précis par localisation et gestion des permissions géolocalisation."),
]

for i, (icon, title, body) in enumerate(challenges):
    col = i % 2
    row = i // 2
    lx = 0.7 + col * 6.3
    ty = 1.85 + row * 2.6
    add_rect(sl, lx, ty, 6.0, 2.35, fill=C_CARD, line_color=RGBColor(0x2A,0x3A,0x4A), line_w=Pt(1))
    # left accent bar
    add_rect(sl, lx, ty, 0.07, 2.35, fill=C_GREEN)
    add_text(sl, icon,  lx+0.25, ty+0.18, 0.7,  0.55, size=22)
    add_text(sl, title, lx+1.1,  ty+0.18, 4.7,  0.42, size=14, bold=True, color=C_GREEN)
    add_text(sl, body,  lx+1.1,  ty+0.72, 4.7,  1.4,  size=10, color=C_GRAY, wrap=True)

# ══════════════════════════════════════════════════════════════
# SLIDE 9 — CONCLUSION
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
bg(sl, C_BG_DARK, RGBColor(0x0A,0x1A,0x12))
add_rect(sl, 0, 0, 0.12, 7.5, fill=C_GREEN)
deco_circle(sl, 9.5, -1, 4.5)
deco_circle(sl, -1, 4.5, 3.5)

section_label(sl, "06 — CONCLUSION & PERSPECTIVES")
slide_title(sl, "Conclusion")
accent_bar(sl, l=0.7, t=1.65)

add_text(sl, "🌙", 5.7, 1.7, 2.0, 1.5, size=54, align=PP_ALIGN.CENTER)
add_text(sl,
    "Curan est une application complète alliant sécurité, spiritualité et technologie moderne.\n"
    "Développée avec Flutter et Firebase, elle offre une expérience utilisateur fluide et personnalisée.",
    1.5, 3.1, 10.3, 0.9, size=13, color=C_GRAY, align=PP_ALIGN.CENTER, wrap=True)

# 3 bottom cards
conc_cards = [
    ("✅", "Réalisé",      "Lecteur, Auth, Biométrie,\nFavoris, Prières, Stats"),
    ("🔭", "Perspectives", "Mode hors-ligne, Notifications,\nTafsir intégré, Widget"),
    ("🙏", "Merci !",      "[ Votre Nom Complet ]\nDes questions ?"),
]
for i, (icon, title, body) in enumerate(conc_cards):
    lx = 0.7 + i * 4.2
    ty = 4.2
    add_rect(sl, lx, ty, 3.9, 2.8, fill=RGBColor(0x0A,0x28,0x1A), line_color=C_GREEN, line_w=Pt(1))
    add_text(sl, icon,  lx+0.2, ty+0.2,  0.7,  0.55, size=22)
    add_text(sl, title, lx+0.2, ty+0.82, 3.5,  0.4,  size=14, bold=True, color=C_GREEN)
    add_text(sl, body,  lx+0.2, ty+1.28, 3.5,  1.3,  size=10, color=C_GRAY, wrap=True)

# ── Save ───────────────────────────────────────────────────────
out = r"c:\Users\Administrator\Documents\GitHub\Quran-app-flutter\Curan_Presentation.pptx"
prs.save(out)
print("Fichier cree : " + out)
